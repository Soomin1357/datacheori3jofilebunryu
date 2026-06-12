import os
import re
import sys
import time
import json
import zlib
import struct
import hashlib
import tempfile
import subprocess
import io
import zipfile
import shutil
import xml.etree.ElementTree as ET
from collections import defaultdict

import numpy as np
import fitz  # PyMuPDF 라이브러리로 구동됨
import olefile
import docx
import pptx
import openpyxl
from PIL import Image
from kiwipiepy import Kiwi
from sentence_transformers import SentenceTransformer
import streamlit as st
import pandas as pd


# NLTK 및 구형 엑셀 전용 라이브러리(xlrd) 패키지 다운로드
@st.cache_resource
def load_nltk_and_words():
    subprocess.run([sys.executable, '-m', 'pip', 'install', 'nltk', 'xlrd', '-q'], check=True)
    import nltk
    nltk.download('words', quiet=True)
    from nltk.corpus import words as en_words
    return set(w.lower() for w in en_words.words())


EN_WORDS = load_nltk_and_words()
import xlrd  # 설치 완료 후 임포트

ABBREV_KEEP = {
    'ai', 'ml', 'dl', 'ocr', 'api', 'url', 'pdf', 'ppt', 'hwp', 'hwpx',
    'doc', 'xls', 'xlsx', 'csv', 'iot', 'ui', 'ux', 'db', 'sql',
    'pc', 'os', 'id', 'ip', 'it', 'bm', 'rd', 'pm', 'pr', 'ir',
    'saas', 'erp', 'crm', 'ceo', 'cto', 'cfo', 'ver', 'no', 'vs', 'idx',
}

# 로컬 PC와 외부 리눅스 서버 환경 자동 감지 및 분기
HAS_OCR = False
try:
    import pytesseract
    
    # 1. 로컬 PC 환경인 경우
    _MYENV_BIN = r"C:\Users\maser\anaconda3\envs\myenv\Library\bin"
    if os.path.exists(_MYENV_BIN):
        _BASE_BIN = r"C:\Users\maser\anaconda3\Library\bin"
        for _p in [_BASE_BIN, _MYENV_BIN]:
            if _p not in os.environ.get("PATH", ""):
                os.environ["PATH"] = _p + ";" + os.environ["PATH"]
        os.environ.setdefault("TESSDATA_PREFIX", r"C:\Users\maser\anaconda3\envs\myenv\share")
        pytesseract.pytesseract.tesseract_cmd = os.path.join(_MYENV_BIN, "tesseract.exe")
        HAS_OCR = True
    else:
        # 2. 외부 리눅스 서버(Streamlit Cloud) 환경인 경우
        if shutil.which("tesseract"):
            pytesseract.pytesseract.tesseract_cmd = "tesseract"
            if os.path.exists("/usr/share/tesseract-ocr/4.00/tessdata"):
                os.environ["TESSDATA_PREFIX"] = "/usr/share/tesseract-ocr/4.00"
            elif os.path.exists("/usr/share/tesseract-ocr/5/tessdata"):
                os.environ["TESSDATA_PREFIX"] = "/usr/share/tesseract-ocr/5"
            HAS_OCR = True
except Exception:
    HAS_OCR = False


class HashCache:
    def __init__(self, cache_path="hash_cache.json"):
        self.cache_path = cache_path
        self._cache = self._load()
        self._dirty = False

    def _load(self):
        if os.path.exists(self.cache_path):
            try:
                with open(self.cache_path, "r", encoding="utf-8") as f:
                    return json.load(f)
            except (json.JSONDecodeError, ValueError):
                backup_path = self.cache_path + ".bak"
                if os.path.exists(self.cache_path):
                    os.rename(self.cache_path, backup_path)
                return {}
        return {}

    def _save(self):
        dir_name = os.path.dirname(self.cache_path) or "."
        temp_file_name = os.path.join(dir_name, f"tmp_{hashlib.md5(self.cache_path.encode()).hexdigest()}")

        try:
            with open(temp_file_name, "w", encoding="utf-8") as f:
                json.dump(self._cache, f, ensure_ascii=False, indent=2)

            os.replace(temp_file_name, self.cache_path)
            self._dirty = False

        except Exception as e:
            if os.path.exists(temp_file_name):
                os.remove(temp_file_name)
            raise e

    def flush(self):
        if self._dirty:
            self._save()

    def compute_hash(self, file_path: str) -> str:
        stat = os.stat(file_path)
        return f"{os.path.abspath(file_path)}|{stat.st_mtime}|{stat.st_size}"

    def check(self, file_path: str):
        file_hash = self.compute_hash(file_path)
        if file_hash in self._cache:
            return True, file_hash, self._cache[file_hash]
        return False, file_hash, None

    def register(self, file_hash: str, result: dict):
        self._cache[file_hash] = result
        self._dirty = True

    def merge_register(self, file_hash: str, extra: dict):
        existing = self._cache.get(file_hash, {})
        self._cache[file_hash] = {**existing, **extra}
        self._dirty = True


class DocumentExtractor:
    MAX_CHARS = 4500

    def __init__(self, max_pages=2):
        self.max_pages = max_pages

    def extract(self, file_path):
        ext = os.path.splitext(file_path)[-1].lower()
        file_name = os.path.basename(file_path)
        extractors = {
            '.pdf': self._extract_pdf,
            '.hwp': self._extract_hwp,
            '.hwpx': self._extract_hwpx,
            '.docx': self._extract_docx,
            '.pptx': self._extract_pptx,
            '.xlsx': self._extract_xlsx,
            '.xls': self._extract_xls,
            '.csv': self._extract_csv,
        }
        try:
            if ext in ['.ppt', '.doc']:
                return self._fallback_filename(file_name, reason=f"구형 확장자({ext}) 본문 추출 미지원")
            extractor_fn = extractors.get(ext)
            if extractor_fn is None:
                return self._fallback_filename(file_name, reason="지원하지 않는 형식")
            text = extractor_fn(file_path)
            if not text or len(text.strip()) < 10:
                return self._fallback_filename(file_name, reason="추출 텍스트 부족")
            return self._trim_to_budget(text)
        except Exception as e:
            return self._fallback_filename(file_name, reason=f"추출 에러: {e}")

    def _extract_xls(self, file_path):
        parts = [f"[META] 파일명: {os.path.basename(file_path)}", "[META] 추출방식: xlrd 구형 셀 스캔", ""]
        try:
            workbook = xlrd.open_workbook(file_path, on_demand=True)
            sheet_names = workbook.sheetnames()
            parts.append(f"[META] 시트 수: {len(sheet_names)}")

            for sheet_idx, sheet_name in enumerate(sheet_names):
                if sheet_idx >= 2: break
                sheet = workbook.sheet_by_name(sheet_name)
                parts.append(sheet_name)

                for row_idx in range(min(sheet.nrows, 20)):
                    row_values = sheet.row_values(row_idx)
                    cells = [str(cell).strip() for cell in row_values if cell is not None]
                    if any(cells):
                        parts.append("  ".join(cells))
            return "\n".join(parts)
        except Exception as e:
            return self._fallback_filename(os.path.basename(file_path), reason=f"xls 파싱 에러: {e}")

    def _extract_hwpx(self, file_path):
        parts = [f"[META] 파일명: {os.path.basename(file_path)}", "[META] 추출방식: HWPX XML 내재 패턴 파싱", ""]
        extracted_content = []
        try:
            with zipfile.ZipFile(file_path, "r") as z:
                sections = sorted([f for f in z.namelist() if f.startswith("Contents/section") and f.endswith(".xml")])
                for section in sections[:self.max_pages]:
                    xml_data = z.read(section)
                    root = ET.fromstring(xml_data)
                    text_elements = root.iter()
                    for elem in text_elements:
                        if elem.text and elem.text.strip():
                            extracted_content.append(elem.text.strip())
            if extracted_content:
                parts.append(" ".join(extracted_content))
                return "\n".join(parts)
            else:
                return self._fallback_filename(os.path.basename(file_path), reason="HWPX XML 내 텍스트 미탐지")
        except Exception as e:
            return self._fallback_filename(os.path.basename(file_path), reason=f"HWPX 파싱 실패: {e}")

    def _extract_csv(self, file_path):
        with open(file_path, "r", encoding="utf-8", errors="ignore") as f:
            lines = [f.readline().strip() for _ in range(30)]
        parts = [f"[META] 파일명: {os.path.basename(file_path)}", ""]
        parts.extend([l for l in lines if l])
        return "\n".join(parts)

    def _extract_pdf(self, file_path):
        doc = fitz.open(file_path)
        parts = [f"[META] 파일명: {os.path.basename(file_path)}", ""]
        total_text_len = 0
        is_scanned = True
        for i, page in enumerate(doc):
            if i >= self.max_pages: break
            blocks = page.get_text("dict")["blocks"]
            page_parts = []
            for block in blocks:
                if block["type"] != 0: continue
                for line_info in block.get("lines", []):
                    spans = line_info.get("spans", [])
                    line_text = "".join(s["text"] for s in spans).strip()
                    if line_text: page_parts.append(line_text)
            seen = set()
            page_parts_dedup = [l for l in page_parts if not (l in seen or seen.add(l))]
            page_text = "\n".join(page_parts_dedup)
            total_text_len += len(page_text)
            if len(page_text.strip()) > 20: is_scanned = False
            parts.append(page_text)
        doc.close()
        if is_scanned or total_text_len < 20:
            return self._light_ocr_pdf(file_path)
        return "\n".join(parts)

    def _extract_hwp(self, file_path):
        if not olefile.isOleFile(file_path):
            return self._fallback_filename(os.path.basename(file_path), reason="OLE 형식이 아닌 HWP")
        f = olefile.OleFileIO(file_path)
        dirs = f.listdir()
        sections = sorted([d for d in dirs if d[0] == 'BodyText' and d[1].startswith('Section')], key=lambda x: x[1])
        if not sections:
            f.close()
            return self._fallback_filename(os.path.basename(file_path), reason="BodyText 섹션 없음")
        parts = [f"[META] 파일명: {os.path.basename(file_path)}", f"[META] 섹션 수: {len(sections)}", ""]
        encodings = ['utf-16-le', 'euc-kr', 'cp949', 'utf-8']
        for section in sections[:self.max_pages]:
            stream = f.openstream(section)
            data = stream.read()
            try:
                decompressed = zlib.decompress(data, -15)
            except zlib.error:
                decompressed = data
            raws = self._hwp_raw_texts(decompressed, encodings)
            if raws:
                parts.extend(raws)
            else:
                fallback = self._extract_hwp_text_fallback(decompressed, encodings)
                if fallback: parts.append(fallback)
        f.close()
        return "\n".join(parts)

    def _hwp_raw_texts(self, data, encodings):
        texts = []
        pos = 0
        while pos + 4 <= len(data):
            try:
                header = struct.unpack_from('<I', data, pos)[0]
                tag_id = header & 0x3FF
                size = (header >> 20) & 0xFFF
                pos += 4
                if size == 0xFFF:
                    if pos + 4 > len(data): break
                    size = struct.unpack_from('<I', data, pos)[0]
                    pos += 4
                if pos + size > len(data): break
                if tag_id == 67 and size > 0:
                    record_data = data[pos:pos + size]
                    text = self._decode_hwp_para_text(record_data, encodings)
                    if text and len(text.strip()) > 1:
                        texts.append(text.strip())
                pos += size
            except (struct.error, Exception):
                pos += (size if 'size' in locals() and size > 0 else 1)
        return texts

    def _decode_hwp_para_text(self, data, encodings):
        try:
            raw_text = data.decode('utf-16-le', errors='ignore')
            cleaned = ""
            i = 0
            while i < len(raw_text):
                ch = raw_text[i]
                code = ord(ch)
                if code == 0:
                    i += 1
                elif code < 0x20:
                    i += 8 if code in (1, 2, 3, 11, 12, 14, 15, 16, 17, 18, 21, 22, 23) else 1
                elif code in (0x0D, 0x0A):
                    cleaned += "\n"
                    i += 1
                else:
                    cleaned += ch
                    i += 1
            if len(cleaned.strip()) > 1: return cleaned.strip()
        except Exception:
            pass
        for enc in encodings[1:]:
            try:
                text = data.decode(enc, errors='ignore')
                text = re.sub(r'[\x00-\x08\x0b\x0c\x0e-\x1f]', '', text)
                if len(text.strip()) > 1: return text.strip()
            except Exception:
                continue
        return None

    def _extract_hwp_text_fallback(self, data, encodings):
        for enc in encodings:
            try:
                text = data.decode(enc, errors='ignore')
                cleaned = re.sub(r'[^가-힣a-zA-Z0-9\s.,!?·\-_()\[\]/%]', '', text)
                cleaned = re.sub(r'\s+', ' ', cleaned).strip()
                korean_words = re.findall(r'[가-힣]{2,}', cleaned)
                if len(korean_words) >= 3: return cleaned
            except Exception:
                continue
        return None

    def _extract_docx(self, file_path):
        doc = docx.Document(file_path)
        parts = [f"[META] 파일명: {os.path.basename(file_path)}", ""]
        current_len = 0
        for para in doc.paragraphs:
            text = para.text.strip()
            if text:
                remaining_budget = 4500 - current_len
                if remaining_budget <= 0: break
                if len(text) > remaining_budget: text = text[:remaining_budget]
                parts.append(text)
                current_len += len(text)
                if current_len >= 4500: break
        return "\n".join(parts)

    def _extract_pptx(self, file_path):
        prs = pptx.Presentation(file_path)
        parts = [f"[META] 파일명: {os.path.basename(file_path)}", f"[META] 총 슬라이드: {len(prs.slides)}", ""]
        current_len = 0
        for slide in prs.slides:
            if current_len >= 4500: break
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    lines = [line.strip() for line in shape.text.strip().split('\n') if line.strip()]
                    for line in lines:
                        remaining_budget = 4500 - current_len
                        if remaining_budget <= 0: break
                        if len(line) > remaining_budget: line = line[:remaining_budget]
                        parts.append(line)
                        current_len += len(line)
                        if current_len >= 4500: break
        return "\n".join(parts)

    def _extract_xlsx(self, file_path):
        wb = openpyxl.load_workbook(file_path, data_only=True, read_only=True)
        parts = [f"[META] 파일명: {os.path.basename(file_path)}", f"[META] 시트 수: {len(wb.sheetnames)}", ""]
        for sheet_idx, sheet_name in enumerate(wb.sheetnames):
            if sheet_idx >= 2: break
            sheet = wb[sheet_name]
            parts.append(sheet_name)
            for i, row in enumerate(sheet.iter_rows(values_only=True)):
                if i >= 20: break
                cells = [str(cell) if cell is not None else "" for cell in row]
                if any(c.strip() for c in cells): parts.append("  ".join(cells).strip())
            parts.append("")
        wb.close()
        return "\n".join(parts)

    def _light_ocr_pdf(self, file_path):
        if not HAS_OCR: return self._fallback_filename(os.path.basename(file_path), reason="OCR 미설치 또는 미지원 환경")
        try:
            doc = fitz.open(file_path)
            page = doc[0]
            rect = page.rect
            clip = fitz.Rect(0, 0, rect.width, rect.height * 0.3)
            pix = page.get_pixmap(clip=clip, dpi=150)
            img = Image.frombytes("RGB", [pix.width, pix.height], pix.samples).convert('L')
            text = pytesseract.image_to_string(img, lang='kor+eng', config='--oem 1 --psm 6')
            doc.close()
            if text and len(text.strip()) > 10:
                return f"[META] 파일명: {os.path.basename(file_path)}\n[META] 추출방식: Light-OCR (상단 30%)\n\n{text.strip()}"
            return self._fallback_filename(os.path.basename(file_path), reason="OCR 결과 부족")
        except Exception as e:
            return self._fallback_filename(os.path.basename(file_path), reason=f"OCR 실패: {e}")

    def _fallback_filename(self, file_name, reason=""):
        clean_name = os.path.splitext(file_name)[0]
        clean_name = re.sub(r'[_\-\.\(\)\[\]{}]', ' ', clean_name)
        clean_name = re.sub(r'\s+', ' ', clean_name).strip()
        parts = [f"[META] 파일명: {file_name}", "[META] 추출방식: 파일명 폴백"]
        if reason: parts.append(f"[META] 사유: {reason}")
        parts.extend(["", clean_name])
        return "\n".join(parts)

    def _trim_to_budget(self, text):
        if len(text) <= self.MAX_CHARS: return text
        lines = text.split('\n')
        meta_lines = [l for l in lines if l.startswith('[META]')]
        content_lines = [l for l in lines if not l.startswith('[META]')]
        meta_text = "\n".join(meta_lines) + "\n\n"
        remaining = self.MAX_CHARS - len(meta_text)
        content_text = "\n".join(content_lines)
        if len(content_text) > remaining:
            content_text = content_text[:remaining - 20] + "\n\n[... 입력 예산 초과로 생략]"
        return meta_text + content_text


FILENAME_RULES = {
    "1. 공고_지침_양식": ["공고", "지침", "양식", "안내서", "요청서", "지시서", "안내", "요청", "지시", "공고문", "지침서", "양식서"],
    "2. 사업계획서 수행계획서": ["사업계획서", "수행계획서", "계획서", "창업", "벤처", "제안서", "계획"],
    "3. 조사_참고자료": ["시장조사", "동향", "분석", "참고", "트렌드", "통계", "별표", "조사", "별첨"],
    "4. 중간_최종 결과물 및 보고서": ["결과보고서", "최종보고서", "중간보고서", "완료보고서", "보고서", "성과물", "최종", "보고"],
    "5. 발표자료": ["발표", "발표자료", "프레젠테이션", "슬라이드", "피치덱"],
    "6. 견적_계약_정산": ["견적서", "계약서", "세금계산서", "정산서", "내역서", "영수증", "지출", "견적", "계약", "정산", "세금", "내역", "계산"],
    "7. 기업 인증서": ["등록증", "등본", "확인서", "증명서", "재무제표", "인증서", "증서"],
    "8. 동의서": ["동의서", "서약서", "동의", "서약"]
}

ANCHOR_DESCRIPTIONS = {
    "1. 공고_지침_양식": "정부나 기관에서 발행한 공식 공고문, 모집 공고, 입찰 공고입니다.",
    "2. 사업계획서 수행계획서": "창업, 연구개발, 국책과제, 기술개발 사업에 제출하는 사업계획서 혹은 과제 제안서입니다.",
    "3. 조사_참고자료": "외부 환경, 시장, 기술, 정책을 조사하고 분석한 참고 자료, 통계 및 규정집입니다.",
    "4. 중간_최종 결과물 및 보고서": "사업이나 연구를 수행한 후 성과와 개발 내용을 담아 제출하는 중간/최종 보고서입니다.",
    "5. 발표자료": "발표나 평가 심사를 위해 파워포인트 등으로 제작된 프레젠테이션 피치덱 슬라이드입니다.",
    "6. 견적_계약_정산": "거래 비용 지급 및 계약 내용을 증빙하는 견적서, 계약서, 세금계산서, 정산 지출 내역 관련 서류입니다.",
    "7. 기업 인증서": "사업자등록증, 등본, 각종 확인서 등 기업의 존재 자격을 증명하는 공식 서류입니다.",
    "8. 동의서": "개인정보 수집 동의서 및 확약서 등 서약 위주의 행정 문서입니다."
}

CATEGORIES_LIST = list(FILENAME_RULES.keys()) + ["9. 기타"]
TEXT_LIMIT = 1500


@st.cache_resource
def load_sbert_model():
    model = SentenceTransformer('BM-K/KoSimCSE-roberta')
    anchor_vecs = {cat: model.encode(text, show_progress_bar=False) for cat, text in ANCHOR_DESCRIPTIONS.items()}
    return model, anchor_vecs


def prepare_text_for_embedding(raw_text: str) -> str:
    lines = raw_text.splitlines()
    file_name, content_parts = "", []
    for l in lines:
        if l.startswith("[META] 파일명:"):
            fn = l.replace("[META] 파일명:", "").strip()
            fn = os.path.splitext(fn)[0]
            fn = re.sub(r'[_\-\.\(\)\[\]{}0-9]', ' ', fn)
            file_name = re.sub(r'\s+', ' ', fn).strip()
        elif l.startswith("[META]"):
            continue
        elif l.strip():
            content_parts.append(l.strip())
    combined = ((file_name + " ") + " ".join(content_parts)).strip()
    tokens = combined.split()
    cleaned = []
    for tok in tokens:
        if tok.isalpha() and tok.islower():
            if tok not in EN_WORDS and tok not in ABBREV_KEEP: continue
        cleaned.append(tok)
    combined = ' '.join(cleaned)
    combined = re.sub(r'\s+', ' ', combined).strip()
    return combined[:TEXT_LIMIT]


def cosine_sim(v1, v2):
    return float(np.dot(v1, v2) / (np.linalg.norm(v1) * np.linalg.norm(v2)))


# ==========================================
# Streamlit 레이아웃 및 상태 초기화
# ==========================================
st.set_page_config(page_title="멀티 문서 분류 및 ZIP 패킹 패키지", layout="wide")

if "classified_df" not in st.session_state:
    st.session_state.classified_df = None
if "uploaded_file_bytes" not in st.session_state:
    st.session_state.uploaded_file_bytes = {}
if "zip_ready" not in st.session_state:
    st.session_state.zip_ready = False
if "zip_buffer" not in st.session_state:
    st.session_state.zip_buffer = None

header_col1, header_col2 = st.columns([2.5, 1.5])

with header_col1:
    st.title("📂 다중 문서 타입 자동 분류 및 패킹 시스템")
    st.caption("여러 문서를 업로드하면 AI 모델이 분석하여 카테고리별 정렬 구조의 가상 아카이브를 제공합니다.")

with header_col2:
    st.markdown("<br>", unsafe_allow_html=True)
    with st.expander("⚠️ 시스템 주의 사항 및 사용 안내", expanded=True):
        st.markdown("""
        **[시스템 안내 사항]**
        1. **이 시스템은 AI 모델이기 때문에 오분류가 발생할 수 있습니다.**
        2. **사용자의 컴퓨터 환경(CPU/GPU 사양)에 따라 속도가 오래 걸릴 수 있습니다.**
        3. 에러가 뜰 때에는 새로고침하거나 `일괄 자동 분류 시작`을 한 번 더 누르면 해결됩니다.

        **[사용 프로세스]**
        1. 파일을 일괄 업로드한 뒤 `일괄 자동 분류 시작`을 누릅니다.
        2. 하단 표에서 분류 결과를 확인 및 수정합니다.
        3. `최종 ZIP 아카이브 생성` 버튼을 통해 압축 파일을 다운로드합니다.
        """)

st.divider()

SBERT_MODEL, ANCHOR_VECS = load_sbert_model()

uploaded_files = st.file_uploader(
    "분석할 문서들을 복수 선택하여 업로드 하세요. (.pdf, .hwp, .hwpx, .docx, .pptx, .xlsx, .xls, .csv)",
    type=["pdf", "hwp", "hwpx", "docx", "pptx", "xlsx", "xls", "csv"],
    accept_multiple_files=True
)

if uploaded_files:
    current_names = set(f.name for f in uploaded_files)
    existing_names = set(st.session_state.uploaded_file_bytes.keys())

    if current_names != existing_names:
        st.session_state.classified_df = None
        st.session_state.zip_ready = False
        st.session_state.zip_buffer = None
        st.session_state.uploaded_file_bytes = {f.name: f.getbuffer() for f in uploaded_files}
else:
    st.session_state.classified_df = None
    st.session_state.uploaded_file_bytes = {}
    st.session_state.zip_ready = False
    st.session_state.zip_buffer = None

# ==========================================
# 실행 제어 및 핵심 분석 엔진 구역
# ==========================================
if uploaded_files:
    if st.session_state.classified_df is None:
        st.info(f"현재 {len(uploaded_files)}개의 파일이 대기 중입니다. 아래 버튼을 누르면 인공지능 분석이 시작됩니다.")
        if st.button("일괄 자동 분류 시작", type="primary", use_container_width=True):

            summary_results = []
            cache = HashCache()
            extractor = DocumentExtractor()
            kiwi = Kiwi()

            user_complex_nouns = [
                "공고", "공고문", "안내", "안내서", "요청", "요청서", "지시", "지시서", "지침", "지침서", "양식", "양식서",
                "계획서", "사업계획서", "수행계획서", "제안서", "창업", "벤처",
                "시장조사", "동향", "분석", "참고", "트렌드", "통계", "조사", "별표", "별첨",
                "보고서", "결과보고서", "최종보고서", "중간보고서", "완료보고서", "성과물", "발표", "발표자료", "프레젠테이션", "슬라이드", "피치덱",
                "견적", "견적서", "계약", "계약서", "정산", "정산서", "내역", "내역서", "지출", "영수증", "세금", "세금계산서", "계산",
                "등록증", "등본", "확인서", "증명서", "인증서", "증서", "재무제표"
            ]
            for word in user_complex_nouns:
                kiwi.add_user_word(word, tag="NNP")

            progress_bar = st.progress(0)

            with tempfile.TemporaryDirectory() as temp_dir:
                for idx, (f_name, f_bytes) in enumerate(st.session_state.uploaded_file_bytes.items()):
                    file_path = os.path.join(temp_dir, f_name)
                    with open(file_path, "wb") as f:
                        f.write(f_bytes)

                    fname_only, _ = os.path.splitext(f_name)

                    is_hit, file_hash, cached_data = cache.check(file_path)
                    cached_pred = cached_data.get("pred_category") if cached_data else None

                    best_category = None
                    step1_all_scores = {}
                    extracted_text = ""

                    if is_hit and cached_pred:
                        best_category = cached_pred
                    else:
                        # [Step 1] 파일명 매칭
                        clean_name = re.sub(r'[_\-\.\(\)\[\]\{\}\s+]', ' ', fname_only).strip()
                        nouns = []
                        if clean_name:
                            tokens = kiwi.tokenize(clean_name)
                            nouns = [t.form for t in tokens if t.tag in {"NNG", "NNP"} and len(t.form) > 1]

                        for category, keywords in FILENAME_RULES.items():
                            matched = set(nouns) & set(keywords)
                            if matched:
                                step1_all_scores[category] = len(matched)

                        if step1_all_scores:
                            sorted_cats = sorted(step1_all_scores.items(), key=lambda x: -x[1])
                            best_category = sorted_cats[0][0]

                        # [Step 2] 본문 규칙 매칭
                        if not best_category:
                            extracted_text = extractor.extract(file_path)
                            cache.register(file_hash, {"text": extracted_text})

                            if extracted_text and "추출방식: 파일명 폴백" not in extracted_text:
                                category_scores = {}
                                for category, keywords in FILENAME_RULES.items():
                                    body_score = sum(extracted_text.count(kw) for kw in keywords)
                                    category_scores[category] = body_score

                                sorted_scores = sorted(category_scores.items(), key=lambda x: x[1], reverse=True)
                                top1_cat, top1_score = sorted_scores[0]

                                if top1_score >= 5:
                                    best_category = top1_cat

                        # [Step 3] 임베딩 유사도 매칭 (임계값 0.5 반영)
                        if not best_category:
                            text_ready = prepare_text_for_embedding(extracted_text if extracted_text else clean_name)
                            vec = SBERT_MODEL.encode(text_ready, show_progress_bar=False)
                            scores = {cat: cosine_sim(vec, av) for cat, av in ANCHOR_VECS.items()}

                            sorted_sbert = sorted(scores.items(), key=lambda x: -x[1])
                            sbert_score = sorted_sbert[0][1]
                            sbert_gap = sorted_sbert[0][1] - sorted_sbert[1][1] if len(sorted_sbert) > 1 else 1.0

                            if sbert_score >= 0.5 and sbert_gap >= 0.03:
                                best_category = sorted_sbert[0][0]
                            else:
                                best_category = "9. 기타"

                        if best_category not in CATEGORIES_LIST:
                            best_category = "9. 기타"

                        cache.merge_register(file_hash, {"pred_category": best_category})
                        cache.flush()

                    summary_results.append({
                        "파일명": f_name,
                        "최종 분류 카테고리": best_category
                    })
                    progress_bar.progress((idx + 1) / len(st.session_state.uploaded_file_bytes))

            st.session_state.classified_df = pd.DataFrame(summary_results)
            st.rerun()

    # ==========================================
    # 2단계 & 3단계: 결과 노출 및 리패킹 구역
    # ==========================================
    else:
        if st.session_state.classified_df is not None and not st.session_state.classified_df.empty:
            
            st.subheader("📊 일괄 자동 분류 완료 현황 (미리보기 및 수정)")
            st.caption("💡 AI가 분류한 결과가 모호하다면 아래 표에서 직접 카테고리를 수동 변경할 수 있습니다.")

            # 1. 파일별 카테고리 매칭 편집기
            edited_df = st.data_editor(
                st.session_state.classified_df,
                column_config={
                    "최종 분류 카테고리": st.column_config.SelectboxColumn(
                        "최종 분류 카테고리",
                        help="수정할 타겟 문서 군집 그룹을 정의합니다.",
                        options=CATEGORIES_LIST,
                        required=True,
                    ),
                    "파일명": st.column_config.TextColumn("파일명", disabled=True),
                },
                use_container_width=True,
                num_rows="fixed",
                key="file_classification_editor"
            )

            if not edited_df.equals(st.session_state.classified_df):
                st.session_state.classified_df = edited_df
                st.session_state.zip_ready = False

            st.divider()

            # 2. 폴더명 일괄 설명 커스텀 설정 구역
            st.subheader("📝 다운로드 폴더별 한 줄 설명 커스텀 설정")
            st.caption("💡 압축 파일 내부 폴더명 뒤 괄호안에 들어갈 설명을 설정합니다. 원하는 문구로 자유롭게 편집해 보세요!")

            desc_data = [{"카테고리 폴더": cat, "폴더 한 줄 설명": desc.split("입니다.")[0]} for cat, desc in ANCHOR_DESCRIPTIONS.items()]
            desc_data.append({"카테고리 폴더": "9. 기타", "폴더 한 줄 설명": "분류 규칙에 미치지 못하는 문서"})
            desc_df = pd.DataFrame(desc_data)

            edited_desc_df = st.data_editor(
                desc_df,
                column_config={
                    "카테고리 폴더": st.column_config.TextColumn("카테고리 폴더", disabled=True),
                    "폴더 한 줄 설명": st.column_config.TextColumn("폴더 한 줄 설명(이 부분을 수정하세요)", required=True),
                },
                use_container_width=True,
                num_rows="fixed",
                key="folder_desc_editor"
            )

            current_folder_descriptions = dict(zip(edited_desc_df["카테고리 폴더"], edited_desc_df["폴더 한 줄 설명"]))

            st.divider()

            # 3. 최종 패킹 다운로드 구역 (무결성 파일 구조 방식)
            st.markdown("### 📥 대용량 분류 완료 패키지 다운로드")

            col_btn1, col_btn2 = st.columns([1, 1])

            with col_btn1:
                if st.button("📦 최종 설정 반영하여 ZIP 파일 생성하기", use_container_width=True, type="secondary"):
                    with st.spinner("지정한 폴더 계층 구조와 커스텀 설명에 맞추어 ZIP 구조 빌드 중..."):
                        
                        with tempfile.TemporaryDirectory() as packing_dir:
                            readme_lines = ["=== 다운로드 문서 분류 구성 안내 가이드 ===\n"]
                            cat_to_eng = {cat: f"Folder_{i+1}" for i, cat in enumerate(CATEGORIES_LIST[:-1])}
                            cat_to_eng["9. 기타"] = "Folder_9"

                            for _, row in st.session_state.classified_df.iterrows():
                                f_name = row["파일명"]
                                assigned_cat = row["최종 분류 카테고리"]

                                file_bytes = st.session_state.uploaded_file_bytes.get(f_name)
                                if file_bytes:
                                    folder_name = cat_to_eng.get(assigned_cat, "Folder_9")
                                    target_folder_path = os.path.join(packing_dir, folder_name)
                                    os.makedirs(target_folder_path, exist_ok=True)

                                    with open(os.path.join(target_folder_path, f_name), "wb") as pf:
                                        pf.write(file_bytes)

                            for cat in CATEGORIES_LIST:
                                eng_name = cat_to_eng.get(cat, "Folder_9")
                                user_desc = current_folder_descriptions.get(cat, "설명 없음")
                                readme_lines.append(f"▶ [{eng_name}] 폴더  <--  {cat} ({user_desc.strip()})")

                            with open(os.path.join(packing_dir, "폴더_설명_안내서.txt"), "w", encoding="utf-8") as rf:
                                rf.write("\n".join(readme_lines))

                            archive_base_path = os.path.join(tempfile.gettempdir(), f"archive_{int(time.time())}")
                            shutil.make_archive(archive_base_path, 'zip', packing_dir)
                            
                            with open(archive_base_path + ".zip", "rb") as zf:
                                st.session_state.zip_buffer = zf.read()
                            
                            if os.path.exists(archive_base_path + ".zip"):
                                os.remove(archive_base_path + ".zip")

                        st.session_state.zip_ready = True
                        st.success("압축 파일이 무결하게 정비되었습니다! 우측 다운로드 버튼을 이용하세요.")

            with col_btn2:
                if st.session_state.zip_ready and st.session_state.zip_buffer is not None:
                    st.download_button(
                        label="💾 전체 분류 결과 ZIP 다운로드",
                        data=st.session_state.zip_buffer,
                        file_name=f"classified_documents_{int(time.time())}.zip",
                        mime="application/zip",
                        use_container_width=True
                    )
                else:
                    st.button("💾 다운로드 대기 중 (ZIP 파일 생성을 먼저 눌러주세요)", disabled=True, use_container_width=True)
